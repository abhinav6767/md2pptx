import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import os

masters_dir = r'Code EZ_ Master of Agents _ Files\Slide Master'
for fname in os.listdir(masters_dir):
    if fname.endswith('.pptx'):
        print(f'\n=== {fname} ===')
        prs = Presentation(os.path.join(masters_dir, fname))
        print(f'Slide width: {prs.slide_width}, height: {prs.slide_height}')
        print(f'Slide width (inches): {prs.slide_width/914400:.2f}, height: {prs.slide_height/914400:.2f}')
        
        for i, sm in enumerate(prs.slide_masters):
            print(f'\nSlide Master {i}:')
            for j, sl in enumerate(sm.slide_layouts):
                print(f'  Layout {j}: {sl.name}')
                for ph in sl.placeholders:
                    print(f'    PH {ph.placeholder_format.idx}: type={ph.placeholder_format.type}, name={ph.name}')
                    print(f'      left={ph.left}, top={ph.top}, w={ph.width}, h={ph.height}')
        
        print(f'\nTotal slides: {len(prs.slides)}')
        for k, slide in enumerate(prs.slides):
            print(f'  Slide {k}: layout={slide.slide_layout.name}')
            for shape in slide.shapes:
                tname = str(shape.shape_type).split('.')[-1].split('(')[0]
                txt = ''
                if shape.has_text_frame:
                    txt = shape.text_frame.text[:60].replace('\n',' ')
                print(f'    {tname} "{shape.name}" text="{txt}"')
