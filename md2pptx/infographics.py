"""
Infographic Generator — Creates programmatic infographics using python-pptx shapes.
Generates process flows, timelines, comparisons, and key metric callouts.
"""

from typing import List, Optional, Tuple
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .design import DesignSystem
from .analyzer import InfographicCandidate


class InfographicGenerator:
    """Generates infographic shapes on PPTX slides."""

    def __init__(self, design: DesignSystem):
        self.design = design

    def generate(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple[float, float, float, float]] = None):
        """Generate an infographic on the given slide. bounds=(left, top, width, height)"""
        
        # 1) Priority: Native PPTX Shapes for structured/consultant layouts
        # We only fall back to AI matplotlib for truly complex data-driven visuals
        native_types = {
            "process": self._generate_process_flow,
            "timeline": self._generate_timeline,
            "comparison": self._generate_comparison,
            "metrics": self._generate_metrics,
            "premium_cards": self._generate_premium_cards,
            "gear_process": self._generate_gear_process,
            "numbered_list": self._generate_numbered_list,
            "swimlane": self._generate_swimlane,
            "icon_cards": self._generate_icon_cards,
            "data_table": self._generate_data_table,
            "vertical_timeline": self._generate_vertical_timeline,
        }
        
        if candidate.infographic_type in native_types:
            native_types[candidate.infographic_type](slide, candidate, bounds)
            return

        # 2) For anything else, attempt AI matplotlib infographic
        try:
            from .agents.infographic_agent import InfographicAgent
            import tempfile
            import os

            out_dir = tempfile.mkdtemp(prefix="md2pptx_infographics_")
            filepath = os.path.join(out_dir, "info_ai.png")
            
            print(f"  [InfographicGenerator] Summoning AI Infographic Agent for stunning {candidate.infographic_type}...")
            agent = InfographicAgent()
            code = agent.generate_infographic_code(
                candidate, 
                filepath,
                theme_colors=self.design.colors.chart_colors()
            )
            
            exec_globals = {"__builtins__": __builtins__, "filepath": filepath}
            exec(code, exec_globals)
            
            if os.path.exists(filepath):
                if bounds:
                    left, top, width, height = [int(x) for x in bounds]
                else:
                    left, top, width, height = int(Inches(0.5)), int(Inches(2.0)), int(Inches(12.33)), int(Inches(5.0))
                slide.shapes.add_picture(filepath, left, top, width=width, height=height)
                return
                
        except Exception as ai_e:
            print(f"  [Warning] AI Infographic generation failed ({ai_e}), falling back to native PPTX shapes.")

        # 3) Final fallback to process flow
        self._generate_process_flow(slide, candidate, bounds)

    def _generate_process_flow(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """Generate a horizontal process flow infographic."""
        items = candidate.items
        if not items:
            return

        n = len(items)
        colors = self.design.colors

        # Layout calculations
        start_left = bounds[0] if bounds else Inches(0.7)
        top = bounds[1] if bounds else Inches(2.2)
        total_width = bounds[2] if bounds else Inches(11.9)
        box_height = bounds[3] if bounds else Inches(2.8)

        # Calculate individual box width and gap
        gap = Inches(0.15)
        available_width = total_width - gap * (n - 1)
        box_width = available_width / n

        # Color cycle for boxes
        box_colors = [
            colors.primary, colors.secondary, colors.accent1,
            colors.accent2, colors.accent3, colors.accent4,
        ]

        for i, item in enumerate(items):
            left = start_left + (box_width + gap) * i
            color = box_colors[i % len(box_colors)]

            # Main box with rounded corners
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(left), int(top), int(box_width), int(box_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = colors.background  # subtle border
            
            try:
                shape.shadow.inherit = False
                shape.shadow.visible = True
            except:
                pass

            # Step number circle
            num_size = Inches(0.5)
            num_left = left + (box_width - num_size) / 2
            num_top = top + Inches(0.2)
            num_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                int(num_left), int(num_top), int(num_size), int(num_size)
            )
            num_shape.fill.solid()
            num_shape.fill.fore_color.rgb = colors.text_light
            num_shape.line.fill.background()

            # Number text
            tf = num_shape.text_frame
            tf.text = str(i + 1)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = color
            tf.word_wrap = True

            # Item text
            text_top = top + Inches(1.0)
            text_height = box_height - Inches(1.2)
            text_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(left + Inches(0.1)), int(text_top),
                int(box_width - Inches(0.2)), int(text_height)
            )
            text_shape.fill.background()
            text_shape.line.fill.background()

            tf = text_shape.text_frame
            tf.word_wrap = True
            # Truncate text for display
            display_text = item[:120] + "..." if len(item) > 120 else item
            tf.text = display_text
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                para.font.size = Pt(14)
                para.font.color.rgb = colors.text_light

        # Add connecting arrows between boxes
        arrow_top = top + box_height / 2
        for i in range(n - 1):
            arrow_left = start_left + box_width * (i + 1) + gap * i
            arrow_shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                int(arrow_left - Inches(0.05)), int(arrow_top - Inches(0.1)),
                int(gap + Inches(0.1)), int(Inches(0.2))
            )
            arrow_shape.fill.solid()
            arrow_shape.fill.fore_color.rgb = colors.text_muted
            arrow_shape.line.fill.background()

    def _generate_timeline(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """Generate a horizontal timeline infographic."""
        items = candidate.items
        if not items:
            return

        n = len(items)
        colors = self.design.colors

        # Center line
        start_left = bounds[0] if bounds else Inches(1.0)
        line_top = bounds[1] + (bounds[3]/2) if bounds else Inches(4.0)
        total_width = bounds[2] if bounds else Inches(11.3)
        line_width = total_width
        
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(start_left), int(line_top),
            int(line_width), int(Inches(0.06))
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = colors.primary
        line_shape.line.fill.background()

        # Timeline points and labels
        spacing = line_width / max(n - 1, 1) if n > 1 else line_width / 2
        box_colors = [colors.primary, colors.secondary, colors.accent1,
                     colors.accent2, colors.accent3]

        for i, item in enumerate(items):
            x = start_left + spacing * i if n > 1 else start_left + line_width / 2
            color = box_colors[i % len(box_colors)]

            # Dot on timeline
            dot_size = Inches(0.25)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                int(x - dot_size / 2), int(line_top - dot_size / 2 + Inches(0.03)),
                int(dot_size), int(dot_size)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

            # Alternating above/below
            if i % 2 == 0:
                text_top = line_top - Inches(1.8)
            else:
                text_top = line_top + Inches(0.6)

            text_width = Inches(2.2)
            text_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(x - text_width / 2), int(text_top),
                int(text_width), int(Inches(1.5))
            )
            text_shape.fill.solid()
            text_shape.fill.fore_color.rgb = color
            text_shape.line.color.rgb = colors.background
            
            try:
                text_shape.shadow.inherit = False
                text_shape.shadow.visible = True
            except:
                pass

            tf = text_shape.text_frame
            tf.word_wrap = True
            display_text = item[:80] + "..." if len(item) > 80 else item
            tf.text = display_text
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                para.font.size = Pt(13)
                para.font.color.rgb = colors.text_light

    def _generate_comparison(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """Generate a side-by-side comparison infographic."""
        items = candidate.items
        if not items:
            return

        colors = self.design.colors
        n = len(items)

        # Layout: cards in a row
        start_left = bounds[0] if bounds else Inches(0.7)
        top = bounds[1] if bounds else Inches(2.0)
        total_width = bounds[2] if bounds else Inches(11.9)
        card_height = bounds[3] if bounds else Inches(4.2)
        gap = Inches(0.2)
        
        available = total_width - gap * (n - 1)
        card_width = available / n
        
        card_colors = [colors.primary, colors.secondary, colors.accent1,
                      colors.accent2, colors.accent3, colors.accent4]
        
        for i, item in enumerate(items):
            left = start_left + (card_width + gap) * i
            color = card_colors[i % len(card_colors)]
            
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(left), int(top), int(card_width), int(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = color
            card.line.color.rgb = colors.background
            
            try:
                card.shadow.inherit = False
                card.shadow.visible = True
            except:
                pass
            
            # Card text
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.2)
            
            display_text = item[:150] + "..." if len(item) > 150 else item
            tf.text = display_text
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                para.font.size = Pt(14)
                para.font.color.rgb = colors.text_light

    def _generate_metrics(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """Generate key metric callout boxes."""
        items = candidate.items
        values = candidate.values
        if not items:
            return

        colors = self.design.colors
        n = min(len(items), 4)

        start_left = bounds[0] if bounds else Inches(1.0)
        top = bounds[1] if bounds else Inches(2.5)
        total_width = bounds[2] if bounds else Inches(11.3)
        box_height = bounds[3] if bounds else Inches(3.0)
        gap = Inches(0.3)

        available = total_width - gap * (n - 1)
        box_width = available / n

        metric_colors = [colors.primary, colors.accent1, colors.accent2, colors.accent3]

        for i in range(n):
            left = start_left + (box_width + gap) * i
            color = metric_colors[i % len(metric_colors)]

            # Metric box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(left), int(top), int(box_width), int(box_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = colors.background
            
            try:
                box.shadow.inherit = False
                box.shadow.visible = True
            except:
                pass

            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)

            # Value (large)
            if i < len(values) and values[i]:
                p = tf.paragraphs[0]
                p.text = values[i]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = colors.text_light

            # Label
            if i < len(items):
                p = tf.add_paragraph()
                p.text = items[i][:70]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(16)
                p.font.color.rgb = colors.text_light

    def _generate_premium_cards(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """Generate premium shadow cards with thin bottom borders for the Hero Header."""
        items = candidate.items
        if not items:
            return

        n = min(len(items), 4)
        start_left = bounds[0] if bounds else Inches(1.0)
        top = bounds[1] if bounds else Inches(5.0)
        total_width = bounds[2] if bounds else Inches(11.33)
        box_height = bounds[3] if bounds else Inches(2.2)
        gap = Inches(0.4)

        available = total_width - gap * (n - 1)
        box_width = available / max(n, 1)

        icons = ["🎯", "💡", "🚀", "📊", "⚙️", "🌍"]

        for i, item in enumerate(items[:n]):
            left = start_left + (box_width + gap) * i
            
            # Shadowed Body
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top), int(box_width), int(box_height))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            box.line.color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            try:
                box.shadow.inherit = False
                box.shadow.visible = True
            except: pass

            # Thin Red Bottom Border (Styling)
            bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left + Inches(0.2)), int(top + box_height - Inches(0.08)), int(box_width - Inches(0.4)), int(Inches(0.08)))
            bottom_line.fill.solid()
            bottom_line.fill.fore_color.rgb = self.design.colors.primary
            bottom_line.line.fill.background()
            
            tf = box.text_frame
            tf.word_wrap = True
            # Icon
            p = tf.paragraphs[0]
            p.text = icons[i % len(icons)] + "\n"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(28)
            
            # Text
            p2 = tf.add_paragraph()
            p2.text = item[:90]
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = Pt(14)
            p2.font.color.rgb = self.design.colors.text_dark

    def _generate_gear_process(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """Generate a vertical/horizontal connected gear layout for the Sidebar layout."""
        items = candidate.items
        if not items:
            return

        n = min(len(items), 3)
        start_left = bounds[0] if bounds else Inches(5.0)
        top = bounds[1] if bounds else Inches(1.5)
        total_width = bounds[2] if bounds else Inches(8.0)
        total_height = bounds[3] if bounds else Inches(5.5)
        
        box_width = total_width / n
        
        for i, item in enumerate(items[:n]):
            left = start_left + box_width * i
            centerY = top + (total_height / 3)
            
            # Add Gear Shape
            gear_size = Inches(1.2)
            gear = slide.shapes.add_shape(
                MSO_SHAPE.GEAR_6, 
                int(left + (box_width/2) - (gear_size/2)), 
                int(centerY - gear_size/2), 
                int(gear_size), int(gear_size)
            )
            gear.fill.solid()
            gear.fill.fore_color.rgb = self.design.colors.text_light
            gear.line.color.rgb = self.design.colors.primary

            # Number strictly inside Gear
            tf = gear.text_frame
            tf.text = f"0{i+1}"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.color.rgb = self.design.colors.primary
            tf.paragraphs[0].font.bold = True
            
            # Connective dashed line (if not last)
            if i < n - 1:
                line_left = left + (box_width/2) + (gear_size/2)
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(line_left), int(centerY), int(box_width - gear_size), int(Inches(0.04)))
                line.fill.solid()
                line.fill.fore_color.rgb = self.design.colors.text_muted
                line.line.fill.background()
            
            # Text Box Below
            tb = slide.shapes.add_textbox(int(left + Inches(0.1)), int(centerY + gear_size/2 + Inches(0.2)), int(box_width - Inches(0.2)), Inches(2.0))
            tb.text_frame.word_wrap = True
            tb.text_frame.text = item[:100]
            self._style_text_frame(tb.text_frame, Pt(14), self.design.colors.text_dark, alignment=PP_ALIGN.CENTER)

    # -------------------------------------------------------------------------
    # PREMIUM CONSULTANT LAYOUTS (Mirroring MBB/Big4 deck aesthetics)
    # -------------------------------------------------------------------------

    def _generate_numbered_list(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """
        Generates a vertical numbered list with large bold step numbers on the left
        and content text on the right. Replicates the Slide 4 'bubble mechanism' style
        from the AI Bubble sample deck.
        """
        items = candidate.items
        if not items:
            return

        n = min(len(items), 6)
        colors = self.design.colors

        left = bounds[0] if bounds else Inches(0.5)
        top = bounds[1] if bounds else Inches(1.8)
        width = bounds[2] if bounds else Inches(12.3)
        total_height = bounds[3] if bounds else Inches(5.2)

        row_height = total_height / n
        num_col_w = Inches(1.0)
        text_col_w = width - num_col_w - Inches(0.2)
        text_col_left = left + num_col_w + Inches(0.2)

        # Color cycle for step numbers
        num_colors = [colors.primary, colors.secondary, colors.accent1, colors.accent2, colors.accent3, colors.accent4]

        for i, item in enumerate(items[:n]):
            row_top = top + row_height * i
            color = num_colors[i % len(num_colors)]

            # --- Large numbered circle ---
            circle_size = min(row_height * 0.75, Inches(0.7))
            circle_left = left + (num_col_w - circle_size) / 2
            circle_top = row_top + (row_height - circle_size) / 2
            
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(circle_left), int(circle_top), int(circle_size), int(circle_size))
            circ.fill.solid()
            circ.fill.fore_color.rgb = color
            circ.line.fill.background()

            tf = circ.text_frame
            tf.text = f"{i+1:02d}"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = colors.text_light

            # --- Separator line (thin horizontal rule between rows) ---
            if i < n - 1:
                sep = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(text_col_left), int(row_top + row_height - Inches(0.02)),
                    int(text_col_w), int(Inches(0.02))
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
                sep.line.fill.background()

            # --- Content text box ---
            txt = slide.shapes.add_textbox(int(text_col_left), int(row_top + Inches(0.1)), int(text_col_w), int(row_height - Inches(0.15)))
            tf2 = txt.text_frame
            tf2.word_wrap = True
            
            # Split item by first sentence break for bold header + detail
            parts = item.split(':', 1)
            if len(parts) == 2:
                p = tf2.paragraphs[0]
                p.text = parts[0].strip()
                p.font.bold = True
                p.font.size = Pt(13)
                p.font.color.rgb = colors.text_dark
                p2 = tf2.add_paragraph()
                p2.text = parts[1].strip()[:150]
                p2.font.size = Pt(12)
                p2.font.color.rgb = colors.text_muted
            else:
                p = tf2.paragraphs[0]
                p.text = item[:180]
                p.font.size = Pt(13)
                p.font.color.rgb = colors.text_dark

    def _generate_swimlane(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """
        Generates a horizontal swimlane/column layout with numbered top circles and
        content cards below — replicating the consulting '01-02-03' column style
        seen in Slides 7-8 of the AI Bubble sample.
        """
        items = candidate.items
        if not items:
            return

        n = min(len(items), 5)
        colors = self.design.colors

        left = bounds[0] if bounds else Inches(0.5)
        top = bounds[1] if bounds else Inches(1.8)
        total_width = bounds[2] if bounds else Inches(12.3)
        total_height = bounds[3] if bounds else Inches(5.2)

        gap = Inches(0.25)
        col_width = (total_width - gap * (n - 1)) / n
        num_colors = [colors.primary, colors.secondary, colors.accent1, colors.accent2, colors.accent3]

        for i, item in enumerate(items[:n]):
            col_left = left + (col_width + gap) * i
            color = num_colors[i % len(num_colors)]

            # --- Number badge at top ---
            badge_size = Inches(0.55)
            badge_left = col_left + (col_width - badge_size) / 2
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(badge_left), int(top), int(badge_size), int(badge_size))
            badge.fill.solid()
            badge.fill.fore_color.rgb = color
            badge.line.fill.background()
            tf = badge.text_frame
            tf.text = f"{i+1:02d}"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = colors.text_light

            # --- Vertical line from badge down to card ---
            vline_top = top + badge_size
            vline_height = Inches(0.3)
            vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(col_left + col_width / 2 - Inches(0.02)), int(vline_top), int(Inches(0.04)), int(vline_height))
            vline.fill.solid()
            vline.fill.fore_color.rgb = color
            vline.line.fill.background()

            # --- Content card ---
            card_top = vline_top + vline_height
            card_height = total_height - badge_size - vline_height
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(col_left), int(card_top), int(col_width), int(card_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xF7, 0xF8, 0xFA)
            card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE8)
            
            # Colored top accent bar on card
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(col_left), int(card_top), int(col_width), int(Inches(0.06)))
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()

            # Text in card
            parts = item.split(':', 1)
            tf2 = card.text_frame
            tf2.word_wrap = True
            tf2.margin_top = Inches(0.2)
            tf2.margin_left = Inches(0.15)
            tf2.margin_right = Inches(0.15)

            if len(parts) == 2:
                p = tf2.paragraphs[0]
                p.text = parts[0].strip()
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.color.rgb = color
                p2 = tf2.add_paragraph()
                p2.text = parts[1].strip()[:200]
                p2.font.size = Pt(10)
                p2.font.color.rgb = colors.text_dark
            else:
                p = tf2.paragraphs[0]
                p.text = item[:200]
                p.font.size = Pt(10)
                p.font.color.rgb = colors.text_dark

    def _generate_icon_cards(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """
        Generates rounded-corner cards with primary-color left accent bars and
        bold title + description text — replicating the Slide 10 'Safeguards' style.
        Works for any 'category: description' content structure.
        """
        items = candidate.items
        if not items:
            return

        n = min(len(items), 6)
        colors = self.design.colors

        left = bounds[0] if bounds else Inches(0.5)
        top = bounds[1] if bounds else Inches(1.8)
        total_width = bounds[2] if bounds else Inches(12.3)
        total_height = bounds[3] if bounds else Inches(5.2)

        # Determine grid: 2 columns if > 3 items
        cols = 2 if n > 3 else 1
        rows = -(-n // cols)  # ceiling division
        col_gap = Inches(0.3)
        row_gap = Inches(0.2)
        col_width = (total_width - col_gap * (cols - 1)) / cols
        row_height = (total_height - row_gap * (rows - 1)) / rows

        accent_colors = [colors.primary, colors.secondary, colors.accent1, colors.accent2, colors.accent3, colors.accent4]

        for i, item in enumerate(items[:n]):
            col_idx = i % cols
            row_idx = i // cols
            card_left = left + col_idx * (col_width + col_gap)
            card_top = top + row_idx * (row_height + row_gap)
            color = accent_colors[i % len(accent_colors)]

            # Background card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(card_left), int(card_top), int(col_width), int(row_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            card.line.color.rgb = RGBColor(0xE5, 0xE5, 0xEC)
            try:
                card.shadow.inherit = False
                card.shadow.visible = True
            except:
                pass

            # Left accent bar
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(card_left), int(card_top + Inches(0.1)), int(Inches(0.07)), int(row_height - Inches(0.2)))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()

            # Text (bold header + description)
            txt_left = card_left + Inches(0.2)
            txt_width = col_width - Inches(0.3)
            txt = slide.shapes.add_textbox(int(txt_left), int(card_top + Inches(0.1)), int(txt_width), int(row_height - Inches(0.2)))
            tf = txt.text_frame
            tf.word_wrap = True

            parts = item.split(':', 1)
            if len(parts) == 2:
                p = tf.paragraphs[0]
                p.text = parts[0].strip()
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = color
                p2 = tf.add_paragraph()
                p2.text = parts[1].strip()[:200]
                p2.font.size = Pt(11)
                p2.font.color.rgb = colors.text_dark
            else:
                p = tf.paragraphs[0]
                p.text = item[:200]
                p.font.size = Pt(11)
                p.font.color.rgb = colors.text_dark

    def _generate_data_table(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """
        Generates a styled indicator/data table with a primary-colored header row
        and alternating row shading. Replicates Slide 6 'market scale indicators'.
        Uses candidate.items as rows in 'Label | Value | Year' format (split by '|').
        """
        items = candidate.items
        if not items:
            return

        colors = self.design.colors

        left = bounds[0] if bounds else Inches(0.5)
        top = bounds[1] if bounds else Inches(1.8)
        total_width = bounds[2] if bounds else Inches(12.3)
        total_height = bounds[3] if bounds else Inches(5.2)

        # Parse rows: expect "Indicator | Value | Year | Source" separated by |
        headers = ["Indicator", "Value", "Year", "Source"]
        rows = []
        for item in items[:10]:
            parts = [p.strip() for p in item.split('|')]
            while len(parts) < 4:
                parts.append("")
            rows.append(parts[:4])

        n_rows = len(rows)
        n_cols = 4
        col_widths = [total_width * 0.42, total_width * 0.2, total_width * 0.15, total_width * 0.23]
        row_h = min(total_height / (n_rows + 1), Inches(0.55))  # +1 for header

        # Header row
        x = left
        for ci, (header, cw) in enumerate(zip(headers, col_widths)):
            hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(top), int(cw), int(row_h))
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = colors.primary
            hdr.line.fill.background()
            tf = hdr.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.1)
            p = tf.paragraphs[0]
            p.text = header
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = colors.text_light
            p.alignment = PP_ALIGN.LEFT
            x += cw

        # Data rows
        alt_bg = RGBColor(0xF2, 0xF6, 0xFF)
        for ri, row in enumerate(rows):
            row_top = top + row_h * (ri + 1)
            x = left
            bg_color = alt_bg if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            for ci, (cell, cw) in enumerate(zip(row, col_widths)):
                cell_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(row_top), int(cw), int(row_h))
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = bg_color
                cell_shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE8)
                cell_shape.line.width = Pt(0.5)
                tf = cell_shape.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.1)
                p = tf.paragraphs[0]
                p.text = str(cell)[:80]
                p.font.size = Pt(10) if ci > 0 else Pt(11)
                p.font.bold = ci == 1  # Bold values
                p.font.color.rgb = colors.primary if ci == 1 else colors.text_dark
                p.alignment = PP_ALIGN.LEFT
                x += cw

    def _generate_vertical_timeline(self, slide, candidate: InfographicCandidate, bounds: Optional[Tuple] = None):
        """
        Generates a vertical timeline with a colored left rail, numbered circles,
        and alternating text boxes — suitable for step-by-step processes or historical
        progressions with more than 4 entries.
        """
        items = candidate.items
        if not items:
            return

        n = min(len(items), 6)
        colors = self.design.colors

        left = bounds[0] if bounds else Inches(0.5)
        top = bounds[1] if bounds else Inches(1.8)
        total_width = bounds[2] if bounds else Inches(12.3)
        total_height = bounds[3] if bounds else Inches(5.2)

        rail_x = left + Inches(0.5)
        rail_width = Inches(0.06)
        text_left = rail_x + Inches(0.6)
        text_width = total_width - Inches(1.3)
        row_height = total_height / n

        num_colors = [colors.primary, colors.secondary, colors.accent1, colors.accent2, colors.accent3, colors.accent4]

        # Draw the vertical rail
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(rail_x), int(top + row_height * 0.3), int(rail_width), int(total_height - row_height * 0.3))
        rail.fill.solid()
        rail.fill.fore_color.rgb = colors.primary
        rail.line.fill.background()

        for i, item in enumerate(items[:n]):
            row_top = top + row_height * i
            color = num_colors[i % len(num_colors)]

            # Dot on rail
            dot_size = Inches(0.3)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                int(rail_x - (dot_size - rail_width) / 2),
                int(row_top + (row_height - dot_size) / 2),
                int(dot_size), int(dot_size))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.color.rgb = colors.text_light
            dot.line.width = Pt(2)

            # Text
            txt = slide.shapes.add_textbox(int(text_left), int(row_top + Inches(0.05)), int(text_width), int(row_height - Inches(0.1)))
            tf = txt.text_frame
            tf.word_wrap = True

            parts = item.split(':', 1)
            if len(parts) == 2:
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = parts[0].strip() + ": "
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = color
                run2 = p.add_run()
                run2.text = parts[1].strip()[:200]
                run2.font.size = Pt(11)
                run2.font.color.rgb = colors.text_dark
            else:
                p = tf.paragraphs[0]
                p.text = item[:200]
                p.font.size = Pt(11)
                p.font.color.rgb = colors.text_dark
