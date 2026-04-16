"""
PPTX Renderer — Builds the final .pptx file from an AgentStoryline and generated assets.
Places content into Slide Master layouts with proper formatting.
"""

import os
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

from .design import (
    DesignSystem, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
    SUBTITLE_LEFT, SUBTITLE_TOP, SUBTITLE_WIDTH, SUBTITLE_HEIGHT,
    BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT,
    COL_LEFT_LEFT, COL_LEFT_WIDTH, COL_RIGHT_LEFT, COL_RIGHT_WIDTH,
    COL_TOP, COL_HEIGHT, CHART_LEFT, CHART_TOP, CHART_WIDTH, CHART_HEIGHT,
    FONT_TITLE, FONT_SUBTITLE, FONT_BODY, FONT_BODY_SMALL, FONT_CAPTION,
    FONT_COVER_TITLE, FONT_COVER_SUBTITLE, FONT_DIVIDER_TITLE,
    FONT_TABLE_HEADER, FONT_TABLE_BODY, FONT_SLIDE_NUMBER,
    MAX_BULLETS_PER_SLIDE, LINE_SPACING, PARAGRAPH_SPACING,
    FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH, FOOTER_HEIGHT,
    MAX_CHARS_PER_BULLET,
)
from .template import TemplateLoader, TemplateInfo
from .agents.models import AgentStoryline, AgentSlide
from .charts import ChartGenerator
from .infographics import InfographicGenerator
from .parser import TableData, Document
from .analyzer import ChartCandidate, InfographicCandidate

# Layouts where we fully manage text ourselves — native placeholders must be hidden
_PREMIUM_LAYOUTS = {"hero_header", "sidebar_split"}

# Padding inside every manually-created text box
_TB_MARGIN = Inches(0.1)


class PPTXRenderer:
    """Renders an AgentStoryline into a .pptx file using the Slide Master template."""

    def __init__(self, template_loader: TemplateLoader):
        self.template = template_loader
        self.info = template_loader.info
        self.design = self.info.design
        self.chart_gen = ChartGenerator(self.design)
        self.infographic_gen = InfographicGenerator(self.design)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def render(self, storyline: AgentStoryline, doc: Document, output_path: str) -> str:
        """Render the storyline to a .pptx file."""
        prs = Presentation(self.info.path)

        # Remove template's existing example slides
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

        layout_map = {l.name: l.index for l in self.info.layouts}

        chart_idx = 0
        image_idx = 0
        mask_shapes = [
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MSO_SHAPE.OVAL,
            MSO_SHAPE.SNIP_1_RECTANGLE,
            MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
            MSO_SHAPE.ROUND_1_RECTANGLE,
        ]

        for i, slide_content in enumerate(storyline.slides):
            safe_title = slide_content.title.encode('ascii', 'ignore').decode('ascii')
            print(f"  Rendering slide {i+1}/{len(storyline.slides)}: [{slide_content.layout_name}] {safe_title[:50]}")

            layout_idx = layout_map.get(slide_content.layout_name, 0)

            # --- Renderer-level guard-rail: fix any remaining layout misassignment ---
            # Even after the layout agent, double-check that Cover/Thank-You/Divider
            # templates only land on the slides they are meant for.
            _lname_lower = (slide_content.layout_name or "").lower()
            _is_first = (i == 0)
            _is_last  = (i == len(storyline.slides) - 1)
            _bad_cover    = any(p in _lname_lower for p in ["cover", "0_title", "title company"]) and not _is_first
            _bad_thankyou = any(p in _lname_lower for p in ["thank you", "thank_you", "1_thank"]) and not _is_last
            _bad_divider  = any(p in _lname_lower for p in ["divider", "section", "c_section"]) and \
                            (_is_first or _is_last or bool(slide_content.bullets))
            if _bad_cover or _bad_thankyou or _bad_divider:
                # Find the content layout index to substitute
                _content_patterns = ["content", "title, content", "title and content", "1_e_title", "title only"]
                for _ln, _li in layout_map.items():
                    if any(p in _ln.lower() for p in _content_patterns):
                        layout_idx = _li
                        print(f"  [Renderer] Guard-rail: replaced bad layout '{slide_content.layout_name}' with '{_ln}' for slide {i+1}")
                        break

            layout = self._get_layout(prs, layout_idx)
            slide = prs.slides.add_slide(layout)

            visual = slide_content.recommended_visual
            has_bullets = bool(slide_content.bullets)

            # --- Premium layout detection (BEFORE any placeholder filling) ---
            # These layouts manage their own text; we must not let the native
            # title placeholder render underneath the custom shapes.
            is_premium_layout = visual in _PREMIUM_LAYOUTS

            # --- Protected layout detection (Cover, Thank You, Divider) ---
            layout_name_lower = (slide_content.layout_name or "").lower()
            is_protected_layout = any(p in layout_name_lower for p in [
                "cover", "thank you", "thank_you", "divider", "section"
            ])

            # --- Step 1: Clear all placeholder text for premium layouts so
            #             they do not bleed through custom shapes (FIX #4) ---
            if is_premium_layout:
                for ph in slide.placeholders:
                    try:
                        ph.text_frame.clear()
                    except Exception:
                        pass

            # --- Step 2: Fill standard placeholders (only for non-premium) ---
            title_set = False
            if not is_premium_layout:
                for ph in slide.placeholders:
                    idx = ph.placeholder_format.idx
                    if idx in (0, 10):  # Main title
                        ph.text = slide_content.title
                        self._style_text_frame(
                            ph.text_frame, FONT_TITLE,
                            self.design.colors.text_dark, bold=True
                        )
                        title_set = True
                    elif idx == 11 and slide_content.subtitle:  # Subtitle
                        ph.text = slide_content.subtitle
                        self._style_text_frame(
                            ph.text_frame, FONT_SUBTITLE,
                            self.design.colors.text_dark
                        )

            # --- Step 3: Title/subtitle fallbacks for non-premium layouts ---
            if not title_set and not is_premium_layout:
                self._add_title_shape(slide, slide_content.title)
                title_set = True

            if slide_content.subtitle and not is_premium_layout:
                subtitle_set_already = any(
                    ph.placeholder_format.idx == 11 and ph.text
                    for ph in slide.placeholders
                )
                if not subtitle_set_already:
                    self._add_subtitle_paragraph(
                        slide, slide_content.subtitle, layout_idx
                    )

            # --- Step 4: Resolve body bounds from Slide Master (FIX #2) ---
            body_bounds = self._resolve_body_bounds(layout_idx)

            has_image = (
                not is_protected_layout
                and visual in ("image", "ultra_dense", "hero_header", "sidebar_split")
                and getattr(slide_content, 'image_url', None)
                and os.path.exists(slide_content.image_url)
            )

            # ----------------------------------------------------------------
            # Content rendering — each branch is mutually exclusive
            # ----------------------------------------------------------------

            # 1. Two-column text
            if slide_content.is_two_column and has_bullets:
                mid = len(slide_content.bullets) // 2
                left_b, top_b, w_b, h_b = body_bounds
                half_w = (w_b - Inches(0.3)) / 2
                self._add_column_bullets(
                    slide, slide_content.bullets[:mid],
                    left_b, top_b, half_w, h_b
                )
                self._add_column_bullets(
                    slide, slide_content.bullets[mid:],
                    left_b + half_w + Inches(0.3), top_b, half_w, h_b
                )

            # 2. Single-column text
            elif visual == "text" and has_bullets:
                self._add_bullets(slide, slide_content.bullets, body_bounds)

            # 3. Mixed: text left + image right
            elif visual == "image" and has_bullets and has_image:
                left_b, top_b, w_b, h_b = body_bounds
                col_w = (w_b - Inches(0.3)) / 2
                # Add image FIRST (behind), then text box on top — fixes z-order overlap
                pic = slide.shapes.add_picture(
                    slide_content.image_url,
                    int(left_b + col_w + Inches(0.3)), int(top_b),
                    int(col_w), int(h_b)
                )
                pic.auto_shape_type = mask_shapes[image_idx % len(mask_shapes)]
                image_idx += 1
                # Text rendered AFTER image so it sits above in z-order
                self._add_column_bullets(
                    slide, slide_content.bullets,
                    left_b, top_b, col_w, h_b
                )

            # 4. Image only (full body zone)
            elif visual == "image" and has_image:
                left_b, top_b, w_b, h_b = body_bounds
                pic = slide.shapes.add_picture(
                    slide_content.image_url,
                    int(left_b), int(top_b), int(w_b), int(h_b)
                )
                pic.auto_shape_type = mask_shapes[image_idx % len(mask_shapes)]
                image_idx += 1

            # 5. Chart (with optional text left)
            elif visual == "chart" and slide_content.visual_reference:
                matched_table = None
                for t in doc.all_tables:
                    if t.title == slide_content.visual_reference or t.title in slide_content.visual_reference:
                        matched_table = t
                        break

                if not matched_table:
                    for t in doc.all_tables:
                        if t.has_numerical_data:
                            matched_table = t
                            break

                if matched_table:
                    candidate = ChartCandidate(
                        table=matched_table,
                        chart_type="bar" if len(matched_table.rows) < 10 else "line",
                        section_title=slide_content.title
                    )
                    chart_path = self.chart_gen.generate(candidate, chart_idx)
                    if chart_path and os.path.exists(chart_path):
                        left_b, top_b, w_b, h_b = body_bounds
                        if has_bullets:
                            col_w = (w_b - Inches(0.3)) / 2
                            self._add_column_bullets(
                                slide, slide_content.bullets,
                                left_b, top_b, col_w, h_b
                            )
                            slide.shapes.add_picture(
                                chart_path,
                                int(left_b + col_w + Inches(0.3)), int(top_b),
                                int(col_w), int(h_b)
                            )
                        else:
                            slide.shapes.add_picture(
                                chart_path,
                                int(CHART_LEFT), int(CHART_TOP),
                                int(CHART_WIDTH), int(CHART_HEIGHT)
                            )
                    chart_idx += 1
                else:
                    self._add_bullets(slide, slide_content.bullets, body_bounds)

            # 5.5. Native Consultant Layouts
            elif visual in ("numbered_list", "swimlane", "icon_cards", "data_table", "vertical_timeline"):
                import copy
                info_type_map = {
                    "numbered_list": "numbered_list",
                    "swimlane": "swimlane",
                    "icon_cards": "icon_cards",
                    "data_table": "data_table",
                    "vertical_timeline": "vertical_timeline",
                }
                info_candidate = InfographicCandidate(
                    infographic_type=info_type_map[visual],
                    title=slide_content.title,
                    items=copy.deepcopy(slide_content.bullets),
                    values=[]
                )
                self.infographic_gen.generate(slide, info_candidate, bounds=body_bounds)

            # 5.6. Ultra-Dense
            elif visual == "ultra_dense":
                left_b, top_b, w_b, h_b = body_bounds
                col_w = (w_b - Inches(0.3)) / 2
                self._add_column_bullets(
                    slide, slide_content.bullets,
                    left_b, top_b, col_w, h_b
                )
                img_exists = has_image
                if img_exists:
                    img_h = h_b / 2 - Inches(0.15)
                    # Primary image: top-right quadrant
                    pic = slide.shapes.add_picture(
                        slide_content.image_url,
                        int(left_b + col_w + Inches(0.3)), int(top_b),
                        int(col_w), int(img_h)
                    )
                    pic.auto_shape_type = mask_shapes[image_idx % len(mask_shapes)]
                    image_idx += 1
                    # Secondary image: bottom-right quadrant
                    secondary_url = getattr(slide_content, 'image_url_secondary', None)
                    if secondary_url and os.path.exists(secondary_url):
                        pic2 = slide.shapes.add_picture(
                            secondary_url,
                            int(left_b + col_w + Inches(0.3)), int(top_b + img_h + Inches(0.3)),
                            int(col_w), int(img_h - Inches(0.15))
                        )
                        pic2.auto_shape_type = mask_shapes[(image_idx + 1) % len(mask_shapes)]

                info_top = top_b + h_b / 2 + Inches(0.15) if img_exists else top_b
                info_h = h_b / 2 - Inches(0.15) if img_exists else h_b
                info_bounds = (
                    left_b + col_w + Inches(0.3),
                    info_top, col_w, info_h
                )

                info_type = (slide_content.visual_reference or "process").lower()
                mapped_type = "process"
                if "timeline" in info_type: mapped_type = "timeline"
                elif "comparison" in info_type: mapped_type = "comparison"
                elif "metrics" in info_type: mapped_type = "metrics"

                import copy
                info_candidate = InfographicCandidate(
                    infographic_type=mapped_type,
                    title=slide_content.title,
                    items=copy.deepcopy(slide_content.bullets)
                )
                self.infographic_gen.generate(slide, info_candidate, bounds=info_bounds)

            # 5.7. Premium: Hero Header (FIX #4 — no native placeholder overlap)
            elif visual == "hero_header":
                hero_h = Inches(3.3)
                if has_image:
                    slide.shapes.add_picture(
                        slide_content.image_url,
                        0, 0, Inches(13.33), int(hero_h)
                    )
                    # Secondary image: small inset thumbnail in bottom-right corner
                    secondary_url = getattr(slide_content, 'image_url_secondary', None)
                    if secondary_url and os.path.exists(secondary_url):
                        thumb_w = Inches(2.8)
                        thumb_h = Inches(1.8)
                        slide.shapes.add_picture(
                            secondary_url,
                            int(Inches(13.33) - thumb_w - Inches(0.3)),
                            int(hero_h + Inches(0.2)),
                            int(thumb_w), int(thumb_h)
                        )

                oval_w = Inches(8.0)
                oval_h = Inches(4.5)
                oval_left = (Inches(13.33) - oval_w) / 2
                oval_top = hero_h - (oval_h / 2)

                oval = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    int(oval_left), int(oval_top), int(oval_w), int(oval_h)
                )
                oval.fill.solid()
                oval.fill.fore_color.rgb = self.design.colors.primary
                oval.line.fill.background()

                tx_box = slide.shapes.add_textbox(
                    int(oval_left + Inches(0.5)),
                    int(oval_top + Inches(1.0)),
                    int(oval_w - Inches(1.0)),
                    Inches(1.5)
                )
                tx_box.text_frame.word_wrap = True
                tx_box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                tx_box.text_frame.text = slide_content.title.upper()
                self._style_text_frame(
                    tx_box.text_frame, Pt(36),
                    self.design.colors.text_light, bold=True,
                    alignment=PP_ALIGN.CENTER
                )

                info_candidate = InfographicCandidate(
                    infographic_type="premium_cards",
                    title="",
                    items=slide_content.bullets
                )
                bounds = (
                    Inches(0.5), hero_h + Inches(1.5),
                    Inches(12.33), Inches(2.5)
                )
                self.infographic_gen.generate(slide, info_candidate, bounds=bounds)

            # 5.8. Premium: Sidebar Split (FIX #4 — no native placeholder overlap)
            elif visual == "sidebar_split":
                sb_w = Inches(4.44)
                sb_top_h = Inches(3.75)

                rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, int(sb_w), int(sb_top_h)
                )
                rect.fill.solid()
                rect.fill.fore_color.rgb = self.design.colors.primary
                rect.line.fill.background()

                tb = slide.shapes.add_textbox(
                    Inches(0.3), Inches(0.5),
                    int(sb_w - Inches(0.6)), Inches(2.5)
                )
                tb.text_frame.word_wrap = True
                tb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                p_title = tb.text_frame.paragraphs[0]
                p_title.text = slide_content.title
                p_title.font.size = Pt(26)
                p_title.font.bold = True
                p_title.font.color.rgb = self.design.colors.text_light

                if slide_content.subtitle:
                    p_sub = tb.text_frame.add_paragraph()
                    p_sub.text = slide_content.subtitle[:120]
                    p_sub.font.size = Pt(12)
                    p_sub.font.bold = False
                    p_sub.font.color.rgb = self.design.colors.text_light

                if has_image:
                    img_top = sb_top_h
                    img_h = Inches(3.75)
                    pic = slide.shapes.add_picture(
                        slide_content.image_url,
                        Inches(0.2), int(img_top + Inches(0.2)),
                        int(sb_w - Inches(0.4)), int(img_h - Inches(0.4))
                    )
                    pic.auto_shape_type = MSO_SHAPE.OVAL

                info_type = slide_content.visual_reference or "process"
                mapped_type = "comparison" if "comparison" in info_type.lower() else "gear_process"

                info_candidate = InfographicCandidate(
                    infographic_type=mapped_type,
                    title="",
                    items=slide_content.bullets
                )
                bounds = (
                    sb_w + Inches(0.3), Inches(1.0),
                    Inches(13.33) - sb_w - Inches(0.6), Inches(5.5)
                )
                self.infographic_gen.generate(slide, info_candidate, bounds=bounds)

            # 6. Generic Infographic
            elif visual == "infographic":
                info_type = (slide_content.visual_reference or "process").lower()
                mapped_type = "process"
                if "timeline" in info_type: mapped_type = "timeline"
                elif "comparison" in info_type: mapped_type = "comparison"
                elif "metrics" in info_type: mapped_type = "metrics"

                candidate = InfographicCandidate(
                    infographic_type=mapped_type,
                    title=slide_content.title,
                    items=slide_content.bullets
                )
                self.infographic_gen.generate(slide, candidate, bounds=body_bounds)

            # Slide number
            self._add_slide_number(slide, prs)

        os.makedirs(
            os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True
        )
        prs.save(output_path)
        return output_path

    # ------------------------------------------------------------------
    # Layout helpers
    # ------------------------------------------------------------------

    def _get_layout(self, prs: Presentation, layout_idx: int):
        sm = prs.slide_masters[0]
        if 0 <= layout_idx < len(sm.slide_layouts):
            return sm.slide_layouts[layout_idx]
        return sm.slide_layouts[0]

    def _resolve_body_bounds(self, layout_idx: int) -> Tuple[int, int, int, int]:
        """
        Return (left, top, width, height) for the body content area.
        Prefers real Slide Master placeholder geometry (FIX #2);
        falls back to design.py constants if not available.
        """
        bounds = self.template.get_body_bounds(layout_idx)
        if bounds:
            left, top, width, height = bounds
            # Sanity-check: the body area must not start above TITLE_HEIGHT
            min_top = int(TITLE_TOP + TITLE_HEIGHT + Inches(0.15))
            if top < min_top:
                top = min_top
                height = max(int(Inches(4.5)), height)
            return (left, top, width, height)
        # Fallback to constants
        return (int(BODY_LEFT), int(BODY_TOP), int(BODY_WIDTH), int(BODY_HEIGHT))

    # ------------------------------------------------------------------
    # Text box helpers (FIX #1 — overflow, FIX #2 — template geometry)
    # ------------------------------------------------------------------

    def _add_title_shape(self, slide, title: str):
        """Fallback manual title when no title placeholder exists."""
        txbox = slide.shapes.add_textbox(
            int(TITLE_LEFT), int(TITLE_TOP),
            int(TITLE_WIDTH), int(Inches(0.9))  # taller to accommodate two-line titles
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = _TB_MARGIN
        tf.margin_right = _TB_MARGIN
        tf.text = title
        self._style_text_frame(tf, FONT_TITLE, self.design.colors.text_dark, bold=True)

    def _add_subtitle_paragraph(self, slide, subtitle: str, layout_idx: int = 0):
        """Adds a subtitle/summary paragraph below the title."""
        # Try to derive actual subtitle position from Slide Master
        title_bounds = self.template.get_placeholder_bounds(layout_idx, 0) or \
                       self.template.get_placeholder_bounds(layout_idx, 10)

        if title_bounds:
            sub_top = title_bounds[1] + title_bounds[3] + int(Inches(0.1))
        else:
            sub_top = int(SUBTITLE_TOP)

        sub_height = int(Inches(0.45))
        txbox = slide.shapes.add_textbox(
            int(SUBTITLE_LEFT), sub_top,
            int(SUBTITLE_WIDTH), sub_height
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = _TB_MARGIN
        tf.margin_right = _TB_MARGIN
        # Limit subtitle length to prevent overflow into body area
        tf.text = subtitle[:150]
        self._style_text_frame(tf, Pt(13), self.design.colors.text_muted, italic=True)

    def _add_bullets(
        self, slide, bullets: List[str],
        body_bounds: Tuple[int, int, int, int]
    ):
        """Single-column bullet list. Adapts font size when many bullets."""
        left, top, width, height = body_bounds
        txbox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE          # Never grow outside the box
        tf.margin_left = _TB_MARGIN
        tf.margin_right = _TB_MARGIN
        tf.margin_top = _TB_MARGIN

        # Reduce font for 5+ bullets to avoid overflow (FIX #1)
        font_size = FONT_BODY if len(bullets) <= 4 else FONT_BODY_SMALL

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Hard-cap each bullet at render time
            display = bullet[:MAX_CHARS_PER_BULLET]
            p.text = f"  \u2022  {display}"
            p.font.size = font_size
            p.font.color.rgb = self.design.colors.text_dark
            p.space_after = Pt(6)
            p.space_before = Pt(3)
            p.alignment = PP_ALIGN.LEFT

    def _add_column_bullets(
        self, slide, bullets: List[str],
        left, top, width, height
    ):
        """Column bullet list — always uses smaller font, hard-caps text."""
        txbox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = _TB_MARGIN
        tf.margin_right = _TB_MARGIN
        tf.margin_top = _TB_MARGIN

        # Reduce further for many bullets in a narrow column
        font_size = FONT_BODY_SMALL if len(bullets) <= 5 else FONT_CAPTION

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            display = bullet[:MAX_CHARS_PER_BULLET]
            p.text = f"  \u2022  {display}"
            p.font.size = font_size
            p.font.color.rgb = self.design.colors.text_dark
            p.space_after = Pt(5)
            p.alignment = PP_ALIGN.LEFT

    def _add_slide_number(self, slide, prs: Presentation):
        slide_num = len(prs.slides)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 4:
                ph.text = str(slide_num)
                return
        txbox = slide.shapes.add_textbox(
            int(FOOTER_LEFT), int(FOOTER_TOP),
            int(FOOTER_WIDTH), int(FOOTER_HEIGHT)
        )
        tf = txbox.text_frame
        tf.text = str(slide_num)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        tf.paragraphs[0].font.size = FONT_SLIDE_NUMBER
        tf.paragraphs[0].font.color.rgb = self.design.colors.text_muted

    def _style_text_frame(
        self, text_frame, font_size: Pt, color: RGBColor,
        bold: bool = False, italic: bool = False, alignment=None
    ):
        for para in text_frame.paragraphs:
            para.font.size = font_size
            para.font.color.rgb = color
            para.font.bold = bold
            para.font.italic = italic
            if alignment:
                para.alignment = alignment
