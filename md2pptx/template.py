"""
Template Loader — Loads and analyzes Slide Master .pptx templates.
Extracts layouts, placeholders, theme colors, and fonts.
"""

import os
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .design import ThemeColors, ThemeFonts, DesignSystem


@dataclass
class LayoutInfo:
    """Information about a slide layout."""
    name: str
    index: int
    placeholders: Dict[int, dict] = field(default_factory=dict)


@dataclass
class TemplateInfo:
    """Parsed template information."""
    path: str
    slide_width: int = 0
    slide_height: int = 0
    layouts: List[LayoutInfo] = field(default_factory=list)
    layout_map: Dict[str, int] = field(default_factory=dict)
    design: DesignSystem = field(default_factory=DesignSystem)

    # Mapped layout indices for each slide type
    cover_layout_idx: int = 0
    divider_layout_idx: int = -1
    content_layout_idx: int = -1
    blank_layout_idx: int = -1
    thankyou_layout_idx: int = -1


class TemplateLoader:
    """Loads and analyzes a Slide Master .pptx template."""

    # Known layout name patterns for each slide type
    COVER_PATTERNS = ["cover", "1_cover", "2_cover", "0_title", "title company"]
    DIVIDER_PATTERNS = ["divider", "section", "c_section"]
    CONTENT_PATTERNS = ["title only", "title, subtitle", "1_e_title"]
    BLANK_PATTERNS = ["blank"]
    THANKYOU_PATTERNS = ["thank you", "thank_you", "1_thank"]

    def __init__(self, template_path: str):
        self.template_path = template_path
        self.prs = Presentation(template_path)
        self.info = self._analyze()

    def _analyze(self) -> TemplateInfo:
        """Analyze the template and extract all relevant information."""
        info = TemplateInfo(path=self.template_path)
        info.slide_width = self.prs.slide_width
        info.slide_height = self.prs.slide_height

        # Extract layouts
        for sm in self.prs.slide_masters:
            for idx, layout in enumerate(sm.slide_layouts):
                layout_info = LayoutInfo(name=layout.name, index=idx)

                for ph in layout.placeholders:
                    layout_info.placeholders[ph.placeholder_format.idx] = {
                        "type": str(ph.placeholder_format.type),
                        "name": ph.name,
                        "left": ph.left,
                        "top": ph.top,
                        "width": ph.width,
                        "height": ph.height,
                    }

                info.layouts.append(layout_info)
                info.layout_map[layout.name.lower()] = idx

        # Map layout types
        info.cover_layout_idx = self._find_layout(info, self.COVER_PATTERNS, default=0)
        info.divider_layout_idx = self._find_layout(info, self.DIVIDER_PATTERNS, default=-1)
        info.content_layout_idx = self._find_layout(info, self.CONTENT_PATTERNS, default=-1)
        info.blank_layout_idx = self._find_layout(info, self.BLANK_PATTERNS, default=-1)
        info.thankyou_layout_idx = self._find_layout(info, self.THANKYOU_PATTERNS, default=-1)

        # If no content layout, fallback to blank
        if info.content_layout_idx == -1:
            info.content_layout_idx = info.blank_layout_idx

        # Extract theme
        info.design = self._extract_design()

        return info

    def _find_layout(self, info: TemplateInfo, patterns: List[str], default: int = -1) -> int:
        """Find a layout index matching any of the given name patterns."""
        for name_lower, idx in info.layout_map.items():
            for pattern in patterns:
                if pattern in name_lower:
                    return idx
        return default

    def _extract_design(self) -> DesignSystem:
        """Extract design tokens from the template's theme."""
        colors = ThemeColors()
        fonts = ThemeFonts()

        try:
            sm = self.prs.slide_masters[0]
            theme = sm.element
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            theme_elements = theme.findall(".//a:theme/a:themeElements/a:clrScheme", ns)

            if theme_elements:
                clr = theme_elements[0]
                
                def extract_rgb(elem):
                    if elem is None:
                        return None
                    srgb = elem.find("a:srgbClr", ns)
                    if srgb is not None:
                        val = srgb.get("val", "")
                        if len(val) == 6:
                            return RGBColor(int(val[0:2],16), int(val[2:4],16), int(val[4:6],16))
                    # Also handle sysClr (system color reference)
                    sys_clr = elem.find("a:sysClr", ns)
                    if sys_clr is not None:
                        last_clr = sys_clr.get("lastClr", "")
                        if len(last_clr) == 6:
                            return RGBColor(int(last_clr[0:2],16), int(last_clr[2:4],16), int(last_clr[4:6],16))
                    return None

                # Extract all theme color slots
                dk1_rgb = extract_rgb(clr.find("a:dk1", ns))
                dk2_rgb = extract_rgb(clr.find("a:dk2", ns))
                lt1_rgb = extract_rgb(clr.find("a:lt1", ns))
                ac1_rgb = extract_rgb(clr.find("a:accent1", ns))
                ac2_rgb = extract_rgb(clr.find("a:accent2", ns))
                ac3_rgb = extract_rgb(clr.find("a:accent3", ns))
                ac4_rgb = extract_rgb(clr.find("a:accent4", ns))
                ac5_rgb = extract_rgb(clr.find("a:accent5", ns))
                ac6_rgb = extract_rgb(clr.find("a:accent6", ns))

                # Map to our design system intelligently:
                # dk1 = primary dark (usually the brand color for headers)
                # accent1/2 = supporting colors
                if dk1_rgb: colors.primary = dk1_rgb
                if dk2_rgb: colors.secondary = dk2_rgb
                if lt1_rgb: colors.background = lt1_rgb
                if ac1_rgb: colors.accent1 = ac1_rgb
                if ac2_rgb: colors.accent2 = ac2_rgb
                if ac3_rgb: colors.accent3 = ac3_rgb
                if ac4_rgb: colors.accent4 = ac4_rgb
                if ac5_rgb: colors.accent5 = ac5_rgb
                # If secondary was not set via dk2, use accent1 as secondary
                if dk2_rgb is None and ac1_rgb: colors.secondary = ac1_rgb
                
                # Auto-detect dark vs light text based on primary brightness
                if dk1_rgb:
                    brightness = (dk1_rgb[0]*299 + dk1_rgb[1]*587 + dk1_rgb[2]*114) / 1000
                    if brightness < 128:
                        colors.text_dark = dk1_rgb
                    else:
                        colors.text_dark = RGBColor(0x1A, 0x1A, 0x2E)
                
                # Table header uses primary
                colors.table_header_bg = colors.primary
                colors.divider_bg = colors.primary
        except Exception:
            pass  # Use defaults if theme extraction fails

        # Try extracting fonts from slide master text styles
        try:
            sm = self.prs.slide_masters[0]
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                  "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
            # Look for <p:txStyles> -> <p:titleStyle> or <p:bodyStyle>
            tx_styles = sm.element.find(".//p:txStyles", ns)
            if tx_styles is not None:
                latin = tx_styles.find(".//a:latin", ns)
                if latin is not None and latin.get("typeface"):
                    tf = latin.get("typeface")
                    if tf and not tf.startswith("+"):
                        fonts.heading = tf
                        fonts.body = tf
        except Exception:
            pass

        return DesignSystem(colors=colors, fonts=fonts)

    def get_presentation(self) -> Presentation:
        """Return a fresh Presentation based on the template."""
        return Presentation(self.template_path)

    def get_layout(self, layout_idx: int):
        """Get a specific slide layout by index."""
        sm = self.prs.slide_masters[0]
        if 0 <= layout_idx < len(sm.slide_layouts):
            return sm.slide_layouts[layout_idx]
        return sm.slide_layouts[0]

    def get_placeholder_bounds(
        self, layout_idx: int, ph_idx: int
    ) -> Optional[Tuple[int, int, int, int]]:
        """
        Return (left, top, width, height) in EMU for a placeholder in a layout.
        Returns None if the placeholder is not found in that layout.
        """
        sm = self.prs.slide_masters[0]
        layouts = sm.slide_layouts
        if not (0 <= layout_idx < len(layouts)):
            return None
        layout = layouts[layout_idx]
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == ph_idx:
                return (ph.left, ph.top, ph.width, ph.height)
        return None

    def get_body_bounds(
        self, layout_idx: int
    ) -> Optional[Tuple[int, int, int, int]]:
        """
        Return body placeholder (idx=1) bounds for a layout, falling back to
        the master's body placeholder if the layout doesn't define one.
        """
        # Try layout body (idx=1)
        bounds = self.get_placeholder_bounds(layout_idx, 1)
        if bounds:
            return bounds
        # Try master body placeholder
        sm = self.prs.slide_masters[0]
        for ph in sm.placeholders:
            if ph.placeholder_format.idx == 1:
                return (ph.left, ph.top, ph.width, ph.height)
        return None


def find_templates(templates_dir: str) -> Dict[str, str]:
    """Find all available template .pptx files in a directory."""
    templates = {}
    if os.path.isdir(templates_dir):
        for fname in os.listdir(templates_dir):
            if fname.endswith(".pptx") and fname.startswith("Template_"):
                key = fname.replace("Template_", "").replace(".pptx", "").strip()
                # Create a short key
                short_key = key.split("_")[0].lower() if "_" in key else key[:20].lower()
                templates[short_key] = os.path.join(templates_dir, fname)
    return templates
